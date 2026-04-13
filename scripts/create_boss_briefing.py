from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PAPER_DIR = ROOT / "writing_outputs" / "20260412_autoteakg_silver_paper"
OUT_DIR = PAPER_DIR / "boss_briefing"
FIG_DIR = PAPER_DIR / "figures"

PPTX_OUT = OUT_DIR / "AutoTeaKG_Silver_boss_briefing_cn.pptx"
DOC_OUT = OUT_DIR / "AutoTeaKG_Silver_boss_briefing_cn.md"

COLORS = {
    "bg": RGBColor(247, 243, 234),
    "ink": RGBColor(28, 52, 56),
    "tea": RGBColor(132, 88, 48),
    "sage": RGBColor(92, 130, 104),
    "gold": RGBColor(220, 166, 74),
    "coral": RGBColor(197, 93, 72),
    "white": RGBColor(255, 255, 255),
    "muted": RGBColor(92, 102, 104),
}


SLIDES = [
    {
        "title": "AutoTeaKG-Silver 项目进展汇报",
        "subtitle": "茶功能活性文献的自动化证据图谱与论文投稿准备",
        "type": "title",
        "notes": "这页开场说明：本项目已经从想法、文献调研、数据抓取、LLM 标注、KG 构建、图表生成推进到论文 v6 和演示前端。今天汇报重点是项目价值、已有成果、风险和下一步需要老板决策的事项。",
    },
    {
        "title": "一句话：把茶健康研究从“散文献”变成“证据地图”",
        "type": "message",
        "message": "我们不是证明“茶一定有效”，而是建立一套可追溯、可筛选、可更新的证据基础设施。",
        "points": ["适合支撑综述、课题申报和后续数据库/KG 论文", "核心卖点：自动化、证据分级、不确定性、加工/提取上下文"],
        "notes": "强调项目定位。这里要避免过度声称健康功效，把重点放在证据基础设施。老板需要听到的是：这个方向可以产出数据库/资源型论文，也能支撑后续课题和产品化探索。",
    },
    {
        "title": "为什么值得做：现有茶研究多，但结构化证据少",
        "type": "cards",
        "cards": [
            ("研究很多", "PubMed 中茶活性、加工、微生物、人体证据持续增长"),
            ("证据很碎", "茶类型、成分、模型、终点、加工信息不统一"),
            ("资源缺位", "现有茶数据库偏基因组、转录组或风险物质"),
        ],
        "notes": "这页回答为什么值得做。茶研究本身不新，但把功能活性、证据层级、加工/成分上下文和微生物机制统一进一个可计算图谱，仍然有明显资源空白。",
    },
    {
        "title": "当前核心产物已经成型",
        "type": "numbers",
        "numbers": [("635", "证据记录"), ("1,989", "KG 节点"), ("8,195", "KG 关系"), ("190", "微生物相关记录")],
        "notes": "用数字说明不是停留在概念。当前已有可展示的 KG v3，包括记录、节点、边和微生物机制子集。后面论文和前端都基于这版数据。",
    },
    {
        "title": "技术路线：自动化构建 + 不确定性显式标注",
        "type": "image_points",
        "image": "fig_graphical_abstract.png",
        "points": ["PubMed 检索与去重", "GLM5 自动抽取 evidence records", "PMC methods 补加工/提取信息", "导出 KG v3 + query tables + figures"],
        "notes": "讲技术路线，但不要陷入实现细节。重点是 pipeline 是可复现的，且不是一次性 Excel 表。每条证据保留 PMID、研究类型、证据等级、置信度和不确定性标签。",
    },
    {
        "title": "关键发现 1：证据层级差异很大",
        "type": "image_points",
        "image": "fig_activity_evidence_heatmap.png",
        "points": ["很多功能说法仍以动物研究和综述为主", "人体干预/观察证据集中在少数方向", "证据分级是论文的必要贡献"],
        "notes": "这页说明结果价值。证据热图让我们能看到哪个方向证据更强，哪个方向还偏 preclinical。这是资源论文的重要图。",
    },
    {
        "title": "关键发现 2：加工/提取上下文是最大缺口",
        "type": "image_points",
        "image": "fig_context_coverage.png",
        "points": ["processing context: 146 → 183", "extraction context: 154 → 185", "仍有 303 条缺少关键上下文"],
        "notes": "强调这是一个有意义的负结果。不是模型没做好，而是摘要和开放全文中本来就常常缺少工艺信息。这可以支撑论文中对 full text 和 methods extraction 的讨论。",
    },
    {
        "title": "关键发现 3：图谱能支持机制路径查询",
        "type": "image_points",
        "image": "fig_graph_query_polysaccharide_microbiome.png",
        "points": ["示例：茶多糖 → 菌群/SCFAs → 肝脏炎症/脂质沉积", "保留证据等级和不确定性", "比关键词检索更适合机制整理"],
        "notes": "这是回应 reviewer 会问的“为什么要 KG”。用茶多糖-菌群-代谢物-宿主表型路径作为例子，说明图谱不是单纯可视化，而是能做结构化查询。",
    },
    {
        "title": "论文状态：已到 v6，适合资源/数据库方向",
        "type": "status",
        "items": [
            ("论文草稿", "v6 Database-oriented PDF 已编译"),
            ("前端展示", "静态 KG explorer 已完成"),
            ("公众汇报", "普通观众版 PPT 已完成"),
            ("代码仓库", "已推送 GitHub: AutoTeaKG_Silver"),
        ],
        "notes": "向老板汇报项目管理状态。重点是已有论文、前端、PPT、仓库四类产物。说明现在不是初步探索，而是进入投稿前打磨阶段。",
    },
    {
        "title": "主要风险与应对",
        "type": "risk",
        "risks": [
            ("审稿人质疑 LLM 抽取准确性", "已有 48 条验证样本模板，下一步需人工填写并统计"),
            ("认为只是工程流水线", "已加入不确定性模型、stage-wise QC、graph query case"),
            ("版面与目标期刊格式", "已做 Database-oriented v6，仍需按目标期刊模板精修"),
        ],
        "notes": "这页要实事求是。当前最大风险是没有填写完人工验证表，因此准确率不能报告。其它风险已有部分应对。老板需要知道接下来投入在哪里最有效。",
    },
    {
        "title": "建议下一步决策",
        "type": "timeline",
        "steps": [
            ("第 1 步", "完成 48 条 validation worksheet"),
            ("第 2 步", "选择目标期刊：优先 Database / 资源型期刊"),
            ("第 3 步", "按模板做 v7 camera-ready 草案"),
        ],
        "notes": "提出明确下一步。最优先是完成 validation worksheet，因为这是论文可信度的关键。第二是确定目标期刊。第三才是模板和版面细调。",
    },
    {
        "title": "需要老板支持的事项",
        "type": "ask",
        "asks": [
            "确认论文定位：资源/数据库论文，而不是算法论文",
            "安排 1-2 人完成 48 条抽样验证",
            "确定目标期刊或会议方向",
            "决定是否继续扩展到更大规模全文抽取",
        ],
        "notes": "最后明确需要老板决策。不要只说我们做了什么，要说下一步需要什么资源和判断。特别是验证样本需要人工投入，目标期刊需要策略选择。",
    },
    {
        "title": "结论",
        "type": "takeaway",
        "takeaway": "AutoTeaKG-Silver 已具备论文、图谱、前端和汇报材料四类成果；下一步关键是补齐验证指标并完成投稿版定稿。",
        "notes": "收尾强调项目已经具备成型成果，不是单点实验。下一步只要把验证和投稿格式补齐，就可以进入投稿准备。",
    },
]


def add_text(slide, x, y, w, h, text, size=24, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    run = p.runs[0]
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]
    return box


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["bg"]


def title(slide, text):
    add_text(slide, 0.55, 0.35, 12.0, 0.55, text, 28, True, COLORS["ink"])
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.05), Inches(1.05), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["gold"]
    bar.line.color.rgb = COLORS["gold"]


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def card(slide, x, y, w, h, head, body, color):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = COLORS["white"]
    shp.line.color.rgb = color
    shp.line.width = Pt(2)
    add_text(slide, x + 0.18, y + 0.16, w - 0.36, 0.42, head, 17, True, color)
    add_text(slide, x + 0.18, y + 0.76, w - 0.36, h - 0.9, body, 15, False, COLORS["ink"])


def bullets(slide, x, y, w, h, items, size=20):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS["ink"]
        p.space_after = Pt(8)


def image(slide, fname, x, y, w, h):
    path = FIG_DIR / fname
    if not path.exists():
        path = ROOT / "figures" / "kg_v3" / fname
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def make_deck():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for idx, spec in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        bg(slide)
        if spec["type"] == "title":
            add_text(slide, 0.75, 1.15, 11.8, 0.85, spec["title"], 38, True, COLORS["ink"])
            add_text(slide, 0.78, 2.1, 10.8, 0.55, spec["subtitle"], 21, False, COLORS["tea"])
            image(slide, "fig_graphical_abstract.png", 0.95, 3.2, 11.3, 2.6)
        else:
            title(slide, spec["title"])
            if spec["type"] == "message":
                add_text(slide, 1.0, 1.85, 11.3, 1.45, spec["message"], 28, True, COLORS["ink"], PP_ALIGN.CENTER)
                bullets(slide, 2.0, 4.05, 9.3, 1.4, spec["points"], 20)
            elif spec["type"] == "cards":
                for i, (h, b) in enumerate(spec["cards"]):
                    card(slide, 0.75 + i * 4.15, 2.0, 3.45, 2.75, h, b, [COLORS["tea"], COLORS["sage"], COLORS["coral"]][i])
            elif spec["type"] == "numbers":
                for i, (n, lab) in enumerate(spec["numbers"]):
                    x = 0.65 + i * 3.15
                    card(slide, x, 2.1, 2.75, 2.55, n, lab, [COLORS["tea"], COLORS["sage"], COLORS["gold"], COLORS["coral"]][i])
                    add_text(slide, x + 0.18, 2.35, 2.38, 0.72, n, 38, True, [COLORS["tea"], COLORS["sage"], COLORS["gold"], COLORS["coral"]][i], PP_ALIGN.CENTER)
                    add_text(slide, x + 0.18, 3.25, 2.38, 0.5, lab, 17, True, COLORS["ink"], PP_ALIGN.CENTER)
            elif spec["type"] == "image_points":
                image(slide, spec["image"], 0.7, 1.35, 7.0, 4.8)
                bullets(slide, 8.05, 1.6, 4.5, 4.2, spec["points"], 20)
            elif spec["type"] == "status":
                for i, (h, b) in enumerate(spec["items"]):
                    y = 1.55 + i * 1.1
                    card(slide, 1.0, y, 11.0, 0.82, h, b, COLORS["sage"] if i % 2 else COLORS["tea"])
            elif spec["type"] == "risk":
                for i, (r, m) in enumerate(spec["risks"]):
                    card(slide, 0.85, 1.55 + i * 1.55, 11.6, 1.15, r, "应对：" + m, COLORS["coral"])
            elif spec["type"] == "timeline":
                for i, (h, b) in enumerate(spec["steps"]):
                    card(slide, 1.05 + i * 4.05, 2.25, 3.35, 2.25, h, b, [COLORS["tea"], COLORS["sage"], COLORS["coral"]][i])
            elif spec["type"] == "ask":
                bullets(slide, 1.4, 1.65, 10.8, 4.6, spec["asks"], 24)
            elif spec["type"] == "takeaway":
                add_text(slide, 1.15, 2.0, 11.1, 2.5, spec["takeaway"], 30, True, COLORS["ink"], PP_ALIGN.CENTER)
        add_text(slide, 12.0, 7.05, 0.6, 0.25, f"{idx:02d}", 9, False, COLORS["muted"], PP_ALIGN.RIGHT)
        notes(slide, spec["notes"])
    prs.save(PPTX_OUT)


def make_doc():
    lines = [
        "# AutoTeaKG-Silver 老板版中文介绍",
        "",
        "## 项目一句话",
        "",
        "AutoTeaKG-Silver 是一个自动化、可追溯、带不确定性标注的茶功能活性证据图谱，用来把分散的茶研究文献整理成可以查询、可以复核、可以持续更新的证据地图。",
        "",
        "## 为什么值得做",
        "",
        "- 茶功能活性研究很多，但分散在不同茶类型、成分、加工方式、实验模型和终点指标中。",
        "- 现有茶数据库多偏基因组、转录组或风险物质，缺少围绕功能活性证据的结构化资源。",
        "- 资源型数据库/KG 论文有明确发表空间，也能支撑后续综述、课题、产品证据梳理和实验选题。",
        "",
        "## 当前已经完成的成果",
        "",
        "- KG v3：635 条证据记录、1,989 个节点、8,195 条关系。",
        "- 论文：已完成 Database-oriented v6 草案，可作为资源/数据库论文继续打磨。",
        "- 前端：已完成静态 KG explorer，可展示指标、分布、子图和证据记录筛选。",
        "- PPT：已有普通观众版和老板汇报版。",
        "- GitHub：已推送到 `git@github.com:ShawnXiha/AutoTeaKG_Silver.git`。",
        "",
        "## 当前核心发现",
        "",
        "- 茶功能活性证据层级差异明显，动物实验、综述、人群研究不能混为一谈。",
        "- 加工/提取上下文是最大缺口：processing context 从 146 提升到 183，extraction context 从 154 提升到 185，但仍有 303 条缺少上下文。",
        "- 不确定性不能隐藏，必须作为图谱属性保留。",
        "- 图谱可以支持机制路径查询，例如茶多糖 -> 肠道菌群/SCFAs -> 肝脏炎症或脂质沉积。",
        "",
        "## 主要风险",
        "",
        "- LLM 抽取准确性需要外部验证，目前已准备 48 条分层验证样本，但还需要人工填写。",
        "- 论文需要明确定位为资源/数据库论文，而不是算法论文。",
        "- 投稿前还需要按目标期刊模板进行版面和补充材料整理。",
        "",
        "## 建议老板决策",
        "",
        "1. 确认投稿定位：优先资源/数据库方向。",
        "2. 安排 1-2 人完成 48 条验证样本。",
        "3. 确定目标期刊或会议。",
        "4. 决定是否继续扩大 full-text extraction 和人工验证规模。",
        "",
        "## PPT 文件",
        "",
        f"- `{PPTX_OUT}`",
    ]
    for i, spec in enumerate(SLIDES, start=1):
        lines.append("")
        lines.append(f"### Slide {i}: {spec['title']}")
        lines.append(spec["notes"])
    DOC_OUT.write_text("\n".join(lines), encoding="utf-8")


if __name__ == "__main__":
    make_deck()
    make_doc()
    print(f"Wrote {PPTX_OUT}")
    print(f"Wrote {DOC_OUT}")
