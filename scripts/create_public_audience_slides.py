from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PAPER_DIR = ROOT / "writing_outputs" / "20260412_autoteakg_silver_paper"
OUT_DIR = PAPER_DIR / "slides_public"
FIG_DIR = PAPER_DIR / "figures"
PPTX_OUT = OUT_DIR / "AutoTeaKG_Silver_public_talk.pptx"
SCRIPT_OUT = OUT_DIR / "AutoTeaKG_Silver_public_talk_script.md"


COLORS = {
    "cream": RGBColor(250, 246, 237),
    "tea": RGBColor(143, 100, 53),
    "deep": RGBColor(34, 68, 73),
    "sage": RGBColor(126, 154, 126),
    "gold": RGBColor(224, 172, 79),
    "coral": RGBColor(205, 101, 75),
    "white": RGBColor(255, 255, 255),
    "gray": RGBColor(88, 92, 94),
}


SLIDES = [
    {
        "title": "一杯茶背后，有一片证据森林",
        "subtitle": "AutoTeaKG-Silver：把茶功能活性文献变成可查询的证据地图",
        "kind": "title",
        "notes": "大家每天都可能喝茶，也常听到绿茶抗氧化、黑茶调肠道、茶多酚有健康作用。但这些说法来自很多不同论文，强弱不一。今天我想讲的不是‘茶一定有什么神奇功效’，而是我们如何把这些分散证据整理成一张可以查询、可以追溯、也能标出不确定性的证据地图。",
    },
    {
        "title": "问题：茶研究很多，但很难比较",
        "kind": "three_cards",
        "cards": [
            ("叫法不同", "绿茶、EGCG、茶多酚、发酵茶、茶多糖常被混用"),
            ("证据不同", "细胞、动物、人群、综述常被放在一起讨论"),
            ("工艺不同", "加工、提取、发酵会改变成分和解释"),
        ],
        "notes": "第一个问题是叫法不统一。第二个问题是证据层级不同，动物实验和人群研究不能直接等同。第三个问题是茶不是一个固定对象，加工和提取会改变它的成分。所以普通读者看到的‘茶有益健康’，背后其实是很多不同层级、不同对象、不同条件的证据。",
    },
    {
        "title": "我们的想法：先不要下结论，先整理证据",
        "kind": "pipeline",
        "steps": ["检索文献", "抽取证据", "标注不确定性", "生成图谱"],
        "notes": "我们的核心思路很朴素：先不急着判断茶有没有某个功效，而是把论文中的证据结构化。每条证据都保留来源、研究类型、活性类别、茶成分、加工或提取信息，以及它的不确定性。这样以后讨论某个健康说法时，就能先看证据地图，而不是只看单篇论文标题。",
    },
    {
        "title": "AutoTeaKG-Silver 是什么？",
        "kind": "image_text",
        "image": "fig_graphical_abstract.png",
        "bullets": [
            "自动从 PubMed 和开放全文中抽取茶功能活性证据",
            "每条关系都保留 PMID、证据级别和不确定性",
            "目标是 silver-standard，不假装每条边都已人工验证",
        ],
        "notes": "AutoTeaKG-Silver 是一个自动生成的证据图谱。Silver-standard 的意思是：它不是人工逐条确认的黄金标准，但每条记录都有来源和不确定性标记。它更像一张可持续更新的证据底图，适合帮助研究者发现哪里证据强，哪里证据弱，哪里需要进一步核查。",
    },
    {
        "title": "当前图谱有多大？",
        "kind": "big_numbers",
        "numbers": [
            ("635", "证据记录"),
            ("1,989", "图谱节点"),
            ("8,195", "图谱关系"),
        ],
        "notes": "当前版本包含 635 条证据记录，接近 2000 个节点，8000 多条关系。这里的节点可以是论文、证据记录、茶类型、成分类别、活性类别、证据级别、机制、菌群、代谢物和不确定性标签。",
    },
    {
        "title": "证据不是一样强：这是图谱必须记录的事",
        "kind": "image_text",
        "image": "fig_activity_evidence_heatmap.png",
        "bullets": [
            "有些结论主要来自动物实验",
            "有些结论有人群研究支持",
            "图谱把证据层级显式分开",
        ],
        "notes": "这张热图展示不同活性类别对应的证据层级。普通传播里常把证据混在一起，但图谱会把动物实验、综述、人群观察、随机试验分开。这样我们就能问：某个茶功能说法到底是动物证据多，还是已经有人群证据？",
    },
    {
        "title": "加工和提取信息，是最容易缺失的关键上下文",
        "kind": "image_text",
        "image": "fig_context_coverage.png",
        "bullets": [
            "加工背景从 146 条提升到 183 条",
            "提取背景从 154 条提升到 185 条",
            "仍有 303 条缺少这类上下文",
        ],
        "notes": "茶的加工和提取会影响成分，但论文摘要里经常不写这些细节。我们先用规则抽取，再用大模型抽取摘要，再去开放全文的方法部分找。最后确实补到一些信息，但仍有 303 条记录缺少加工或提取上下文。这是一个重要结果：不是模型再问一遍就能解决，很多信息本来就没有公开在摘要里。",
    },
    {
        "title": "我们不隐藏不确定性，而是把它放进图谱",
        "kind": "image_text",
        "image": "fig_uncertainty_by_activity.png",
        "bullets": [
            "低不确定性：100 条",
            "中等不确定性：470 条",
            "高不确定性：65 条",
        ],
        "notes": "很多知识图谱看起来很确定，但实际上自动抽取会有不确定性。我们的做法是把不确定性作为图谱的一部分：比如低置信度、机制太泛、缺少加工信息、只有动物证据等。用户可以根据需要只看低不确定性证据，或者专门查看高不确定性区域。",
    },
    {
        "title": "图谱能回答什么问题？一个例子",
        "kind": "image_text",
        "image": "fig_graph_query_polysaccharide_microbiome.png",
        "bullets": [
            "问题：茶多糖如何连接肠道菌群和肝脏表型？",
            "结果：找到丁酸、SCFAs、肝脏炎症、脂质沉积等路径",
            "每条路径仍保留证据级别和不确定性",
        ],
        "notes": "这是一个图谱查询例子。我们问：茶多糖如何和肠道菌群、代谢物、肝脏表型连接？图谱能把茶多糖、菌群调节、丁酸或短链脂肪酸、肝脏炎症或脂质沉积放在同一条路径里，并且仍然保留这些证据来自什么研究、是否是动物实验、置信度如何。",
    },
    {
        "title": "这不是“证明茶有效”，而是“让证据可追溯”",
        "kind": "contrast",
        "left": ("不要过度解读", ["不是人工确认的金标准", "不是因果证明", "不是健康建议"]),
        "right": ("真正价值", ["快速看到证据强弱", "发现缺少上下文的区域", "帮助设计下一步实验"]),
        "notes": "这里需要特别强调：这项工作不是在给任何茶产品背书，也不是健康建议。它的价值是把证据变得可追溯、可筛选、可更新。研究者可以用它找证据空白，普通读者也能理解为什么某些说法听起来很强，但实际上证据层级可能还不够。",
    },
    {
        "title": "下一步：从证据地图走向可信知识库",
        "kind": "three_cards",
        "cards": [
            ("人工抽查", "完成 48 条分层样本的字段级验证"),
            ("持续更新", "定期抓取新 PubMed 文献并增量标注"),
            ("图谱应用", "支持机制查询、综述写作和实验选题"),
        ],
        "notes": "下一步我们会做三件事。第一，完成抽样人工验证，量化哪些字段可靠、哪些需要改进。第二，让系统持续更新。第三，把它用于机制查询、综述写作和实验选题。最终目标是让茶功能活性研究从碎片化文本，变成可查询、可复核、可持续更新的证据基础设施。",
    },
    {
        "title": "一句话总结",
        "kind": "takeaway",
        "takeaway": "AutoTeaKG-Silver 把“茶有什么作用？”变成了“哪些证据支持这个说法、证据有多强、还缺什么上下文？”",
        "notes": "最后用一句话总结：我们不是简单回答茶有没有某种功效，而是把问题变得更科学：哪些证据支持这个说法，证据有多强，来自什么研究，还缺少什么上下文。这就是 AutoTeaKG-Silver 的核心价值。谢谢大家。",
    },
    {
        "title": "Q&A",
        "kind": "qa",
        "bullets": [
            "项目：AutoTeaKG-Silver",
            "代码与数据：github.com/ShawnXiha/AutoTeaKG_Silver",
            "欢迎讨论：证据图谱、茶功能活性、自动文献抽取",
        ],
        "notes": "这一页停留在问答环节。可以根据听众问题，补充解释证据级别、银标准图谱、为什么需要人工验证，以及这个系统如何帮助后续茶研究。",
    },
]


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_textbox(slide, x, y, w, h, text, size=28, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Microsoft YaHei"
    if color:
        run.font.color.rgb = color
    return box


def add_title(slide, title):
    add_textbox(slide, 0.55, 0.28, 12.2, 0.65, title, 28, True, COLORS["deep"])
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.98), Inches(1.15), Inches(0.07))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["gold"]
    line.line.color.rgb = COLORS["gold"]


def add_footer(slide, idx):
    add_textbox(slide, 11.85, 7.05, 0.8, 0.25, f"{idx:02d}", 9, False, COLORS["gray"], PP_ALIGN.RIGHT)


def add_notes(slide, notes):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = notes


def add_bullets(slide, x, y, w, h, bullets, size=20):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "Microsoft YaHei"
        p.font.color.rgb = COLORS["deep"]
        p.space_after = Pt(9)
    return box


def add_card(slide, x, y, w, h, heading, body, accent):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = accent
    shape.line.width = Pt(2)
    add_textbox(slide, x + 0.18, y + 0.18, w - 0.36, 0.35, heading, 18, True, accent)
    add_textbox(slide, x + 0.18, y + 0.72, w - 0.36, h - 0.9, body, 15, False, COLORS["deep"])


def add_image(slide, filename, x, y, w, h):
    path = FIG_DIR / filename
    if not path.exists():
        # Fallback to root figure directory if manuscript-local copy is absent.
        path = ROOT / "figures" / "kg_v3" / filename
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def add_pipeline(slide, steps):
    y = 3.0
    xs = [0.8, 3.6, 6.4, 9.2]
    for i, (x, step) in enumerate(zip(xs, steps)):
        add_card(slide, x, y, 2.0, 1.15, f"0{i+1}", step, COLORS["tea"] if i < 2 else COLORS["sage"])
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x + 2.08), Inches(y + 0.35), Inches(0.9), Inches(0.35))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS["gold"]
            arrow.line.color.rgb = COLORS["gold"]


def create_deck():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for idx, spec in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        set_bg(slide, COLORS["cream"])

        if spec["kind"] == "title":
            add_textbox(slide, 0.75, 1.25, 11.8, 1.0, spec["title"], 38, True, COLORS["deep"])
            add_textbox(slide, 0.78, 2.35, 10.8, 0.7, spec["subtitle"], 22, False, COLORS["tea"])
            add_image(slide, "fig_graphical_abstract.png", 0.85, 3.35, 11.6, 2.65)
        else:
            add_title(slide, spec["title"])
            if spec["kind"] == "three_cards":
                for i, (heading, body) in enumerate(spec["cards"]):
                    add_card(slide, 0.75 + i * 4.15, 2.15, 3.45, 2.65, heading, body, [COLORS["tea"], COLORS["sage"], COLORS["coral"]][i])
            elif spec["kind"] == "pipeline":
                add_textbox(slide, 0.8, 1.55, 11.5, 0.55, "把散落的论文，变成可以追溯的证据路径", 23, False, COLORS["tea"], PP_ALIGN.CENTER)
                add_pipeline(slide, spec["steps"])
            elif spec["kind"] == "image_text":
                add_image(slide, spec["image"], 0.75, 1.35, 7.0, 4.75)
                add_bullets(slide, 8.15, 1.65, 4.4, 4.2, spec["bullets"], 20)
            elif spec["kind"] == "big_numbers":
                for i, (num, label) in enumerate(spec["numbers"]):
                    x = 0.9 + i * 4.1
                    add_card(slide, x, 2.05, 3.25, 2.8, num, label, [COLORS["tea"], COLORS["sage"], COLORS["coral"]][i])
                    # overwrite num with larger type for impact
                    add_textbox(slide, x + 0.25, 2.35, 2.75, 0.8, num, 42, True, [COLORS["tea"], COLORS["sage"], COLORS["coral"]][i], PP_ALIGN.CENTER)
                    add_textbox(slide, x + 0.25, 3.35, 2.75, 0.6, label, 22, True, COLORS["deep"], PP_ALIGN.CENTER)
            elif spec["kind"] == "contrast":
                lh, lb = spec["left"]
                rh, rb = spec["right"]
                add_card(slide, 0.9, 1.75, 5.3, 4.35, lh, "\n".join(lb), COLORS["coral"])
                add_card(slide, 7.05, 1.75, 5.3, 4.35, rh, "\n".join(rb), COLORS["sage"])
            elif spec["kind"] == "takeaway":
                add_textbox(slide, 1.15, 2.05, 11.0, 2.6, spec["takeaway"], 30, True, COLORS["deep"], PP_ALIGN.CENTER)
                add_textbox(slide, 1.8, 5.15, 9.7, 0.55, "从“结论”转向“证据路径”", 22, False, COLORS["tea"], PP_ALIGN.CENTER)
            elif spec["kind"] == "qa":
                add_bullets(slide, 1.5, 1.8, 10.2, 3.8, spec["bullets"], 24)

        add_footer(slide, idx)
        add_notes(slide, spec["notes"])

    prs.save(PPTX_OUT)


def write_script():
    lines = ["# AutoTeaKG-Silver 面向普通观众演讲稿", "", "建议时长：12-15 分钟", ""]
    for i, spec in enumerate(SLIDES, start=1):
        lines.append(f"## 第 {i} 页：{spec['title']}")
        lines.append("")
        lines.append(spec["notes"])
        lines.append("")
    SCRIPT_OUT.write_text("\n".join(lines), encoding="utf-8")


if __name__ == "__main__":
    create_deck()
    write_script()
    print(f"Wrote {PPTX_OUT}")
    print(f"Wrote {SCRIPT_OUT}")
